#!/usr/bin/env python3
"""
Professional PowerPoint to PDF Converter for Replit
Converts .pptx/.ppt files to PDF with slide preservation
Compatible with Linux environments
"""

import sys
import json
import traceback
from pathlib import Path
import tempfile
import os

def convert_powerpoint_to_pdf_professional(input_path, output_path):
    """
    Convert PowerPoint presentation to PDF using python-pptx and reportlab
    """
    try:
        from pptx import Presentation
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter, A4, landscape
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.units import inch
        from reportlab.lib import colors
        from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT
        
        # Load PowerPoint presentation
        prs = Presentation(input_path)
        
        # Create PDF document (landscape orientation for slides)
        pdf_doc = SimpleDocTemplate(
            output_path,
            pagesize=landscape(A4),
            rightMargin=36,
            leftMargin=36,
            topMargin=36,
            bottomMargin=36
        )
        
        # Get styles
        styles = getSampleStyleSheet()
        
        # Create custom styles for slides
        title_style = ParagraphStyle(
            'SlideTitle',
            parent=styles['Title'],
            fontSize=18,
            spaceAfter=8,
            textColor=colors.black,
            alignment=TA_CENTER
        )
        
        slide_number_style = ParagraphStyle(
            'SlideNumber',
            parent=styles['Normal'],
            fontSize=12,
            spaceAfter=6,
            textColor=colors.blue,
            alignment=TA_LEFT
        )
        
        content_style = ParagraphStyle(
            'SlideContent',
            parent=styles['Normal'],
            fontSize=11,
            spaceAfter=4,
            textColor=colors.black,
            alignment=TA_LEFT,
            leftIndent=20
        )
        
        bullet_style = ParagraphStyle(
            'SlideBullet',
            parent=styles['Normal'],
            fontSize=10,
            spaceAfter=3,
            textColor=colors.black,
            alignment=TA_LEFT,
            leftIndent=40
        )
        
        # Build content list
        content = []
        
        # Add presentation title
        content.append(Paragraph("PowerPoint to PDF Conversion", title_style))
        content.append(Spacer(1, 20))
        
        # Process each slide
        for slide_idx, slide in enumerate(prs.slides, 1):
            # Add slide number
            content.append(Paragraph(f"Slide {slide_idx}", slide_number_style))
            content.append(Spacer(1, 8))
            
            # Extract text from shapes
            slide_has_content = False
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    slide_has_content = True
                    
                    # Check if it's likely a title (first text or short text)
                    if len(text) < 100 and (slide_idx == 1 or not slide_has_content):
                        content.append(Paragraph(escape_html(text), title_style))
                    else:
                        # Handle bullet points
                        lines = text.split('\n')
                        for line in lines:
                            line = line.strip()
                            if line:
                                if line.startswith('•') or line.startswith('-') or line.startswith('*'):
                                    # Remove bullet marker and add as bullet point
                                    clean_line = line.lstrip('•-* ').strip()
                                    content.append(Paragraph(f"• {escape_html(clean_line)}", bullet_style))
                                else:
                                    content.append(Paragraph(escape_html(line), content_style))
                    
                    content.append(Spacer(1, 4))
            
            # If no text content found, add placeholder
            if not slide_has_content:
                content.append(Paragraph("(Slide contains images or graphics only)", content_style))
            
            # Add spacing between slides
            content.append(Spacer(1, 15))
            
            # Add page break after each slide (except the last one)
            if slide_idx < len(prs.slides):
                content.append(PageBreak())
        
        # Add footer
        footer_style = ParagraphStyle(
            'Footer',
            parent=styles['Normal'],
            fontSize=8,
            textColor=colors.grey,
            alignment=TA_CENTER
        )
        content.append(Spacer(1, 20))
        content.append(Paragraph("Converted by I Love Making PDF", footer_style))
        
        # Build PDF
        pdf_doc.build(content)
        
        print(json.dumps({
            "success": True,
            "message": f"PowerPoint presentation successfully converted to PDF with {len(prs.slides)} slides."
        }))
        return True
        
    except Exception as e:
        print(f"Professional conversion failed: {e}", file=sys.stderr)
        print(traceback.format_exc(), file=sys.stderr)
        return False

def escape_html(text):
    """
    Escape HTML special characters for ReportLab
    """
    if not text:
        return ""
    
    text = text.replace('&', '&amp;')
    text = text.replace('<', '&lt;')
    text = text.replace('>', '&gt;')
    text = text.replace('"', '&quot;')
    text = text.replace("'", '&#x27;')
    
    # Handle special characters
    text = text.replace('\x00', '')
    text = text.replace('\r\n', '<br/>')
    text = text.replace('\n', '<br/>')
    text = text.replace('\r', '<br/>')
    
    return text

def main():
    if len(sys.argv) != 3:
        print(json.dumps({
            "success": False,
            "error": "Usage: python powerpoint_to_pdf_converter.py <input.pptx> <output.pdf>"
        }))
        sys.exit(1)
    
    input_path = sys.argv[1]
    output_path = sys.argv[2]
    
    # Validate input file
    if not os.path.exists(input_path):
        print(json.dumps({
            "success": False,
            "error": "Input file not found"
        }))
        sys.exit(1)
    
    if not (input_path.lower().endswith('.pptx') or input_path.lower().endswith('.ppt')):
        print(json.dumps({
            "success": False,
            "error": "Input file must be a .pptx or .ppt file"
        }))
        sys.exit(1)
    
    try:
        # Convert PowerPoint to PDF
        if convert_powerpoint_to_pdf_professional(input_path, output_path):
            sys.exit(0)
        
        # If conversion fails
        print(json.dumps({
            "success": False,
            "error": "PowerPoint to PDF conversion failed. Please ensure the file is not corrupted."
        }))
        sys.exit(1)
        
    except Exception as e:
        print(json.dumps({
            "success": False,
            "error": f"Conversion failed: {str(e)}"
        }))
        sys.exit(1)

if __name__ == "__main__":
    main()