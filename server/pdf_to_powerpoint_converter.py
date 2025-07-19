#!/usr/bin/env python3
"""
Professional PDF to PowerPoint Converter for Replit
Converts PDF files to .pptx with slide structure
Compatible with Linux environments
"""

import sys
import json
import traceback
from pathlib import Path
import tempfile
import os

def convert_pdf_to_powerpoint_professional(input_path, output_path):
    """
    Convert PDF to PowerPoint using PyMuPDF and python-pptx
    """
    try:
        import fitz  # PyMuPDF
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import io
        from PIL import Image
        
        # Load PDF document
        pdf_doc = fitz.open(input_path)
        
        # Create PowerPoint presentation
        prs = Presentation()
        
        # Process each page as a slide
        for page_num in range(pdf_doc.page_count):
            page = pdf_doc[page_num]
            
            # Add a new slide
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Extract text from the page
            text_blocks = page.get_text("dict")
            
            # Try to extract main content
            slide_title = None
            slide_content = []
            
            if text_blocks and 'blocks' in text_blocks:
                for block in text_blocks['blocks']:
                    if 'lines' in block:
                        for line in block['lines']:
                            if 'spans' in line:
                                line_text = ""
                                for span in line['spans']:
                                    if span.get('text', '').strip():
                                        line_text += span['text'].strip() + " "
                                
                                line_text = line_text.strip()
                                if line_text:
                                    # First non-empty line could be title
                                    if not slide_title and len(line_text) < 100:
                                        slide_title = line_text
                                    else:
                                        slide_content.append(line_text)
            
            # If no structured text, try simple text extraction
            if not slide_title and not slide_content:
                simple_text = page.get_text()
                if simple_text.strip():
                    lines = [line.strip() for line in simple_text.split('\n') if line.strip()]
                    if lines:
                        slide_title = lines[0] if len(lines[0]) < 100 else f"Slide {page_num + 1}"
                        slide_content = lines[1:] if len(lines) > 1 else []
            
            # Set default title if none found
            if not slide_title:
                slide_title = f"Slide {page_num + 1}"
            
            # Add title to slide
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = slide_title
            
            # Format title
            title_paragraph = title_frame.paragraphs[0]
            title_paragraph.font.size = Pt(24)
            title_paragraph.font.bold = True
            title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
            title_paragraph.alignment = PP_ALIGN.CENTER
            
            # Add content to slide
            if slide_content:
                content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
                content_frame = content_box.text_frame
                
                for i, content_line in enumerate(slide_content[:10]):  # Limit to 10 lines
                    if i == 0:
                        content_frame.text = content_line
                        p = content_frame.paragraphs[0]
                    else:
                        p = content_frame.add_paragraph()
                        p.text = content_line
                    
                    # Format content
                    p.font.size = Pt(14)
                    p.font.color.rgb = RGBColor(64, 64, 64)
                    p.alignment = PP_ALIGN.LEFT
            
            # Try to extract and add images (if any)
            try:
                image_list = page.get_images()
                if image_list:
                    # Add first image to slide
                    img_index = image_list[0][0]
                    base_image = pdf_doc.extract_image(img_index)
                    image_bytes = base_image["image"]
                    
                    # Save image temporarily
                    temp_img_path = f"/tmp/slide_image_{page_num}.png"
                    with open(temp_img_path, "wb") as img_file:
                        img_file.write(image_bytes)
                    
                    # Add image to slide
                    slide.shapes.add_picture(temp_img_path, Inches(0.5), Inches(3), Inches(4), Inches(3))
                    
                    # Clean up temp file
                    os.unlink(temp_img_path)
            except Exception as img_error:
                print(f"Could not extract image from page {page_num + 1}: {img_error}", file=sys.stderr)
        
        # Get page count before closing
        total_slides = pdf_doc.page_count
        
        # Save PowerPoint presentation
        prs.save(output_path)
        pdf_doc.close()
        
        print(json.dumps({
            "success": True,
            "message": f"PDF successfully converted to PowerPoint with {total_slides} slides."
        }))
        return True
        
    except Exception as e:
        print(f"Professional conversion failed: {e}", file=sys.stderr)
        print(traceback.format_exc(), file=sys.stderr)
        return False

def main():
    if len(sys.argv) != 3:
        print(json.dumps({
            "success": False,
            "error": "Usage: python pdf_to_powerpoint_converter.py <input.pdf> <output.pptx>"
        }))
        sys.exit(1)
    
    input_path = sys.argv[1]
    output_path = sys.argv[2]
    
    # Validate input file
    if not os.path.exists(input_path):
        print(json.dumps({
            "success": False,
            "error": "Input file not found"
        }))
        sys.exit(1)
    
    if not input_path.lower().endswith('.pdf'):
        print(json.dumps({
            "success": False,
            "error": "Input file must be a .pdf file"
        }))
        sys.exit(1)
    
    try:
        # Convert PDF to PowerPoint
        if convert_pdf_to_powerpoint_professional(input_path, output_path):
            sys.exit(0)
        
        # If conversion fails
        print(json.dumps({
            "success": False,
            "error": "PDF to PowerPoint conversion failed. Please ensure the PDF is not corrupted."
        }))
        sys.exit(1)
        
    except Exception as e:
        print(json.dumps({
            "success": False,
            "error": f"Conversion failed: {str(e)}"
        }))
        sys.exit(1)

if __name__ == "__main__":
    main()